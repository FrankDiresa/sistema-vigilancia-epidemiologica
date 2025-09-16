import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import io
import os

class PowerPointGenerator:
    def __init__(self, data_processor):
        self.data_processor = data_processor
        
        # Mapeo de establecimientos por redes de salud
        self.health_networks = {
            "RED UTCUBAMBA": {
                "establecimientos": [5060, 5044, 7276, 7435, 7006, 5126, 5135, 5136, 5137],
                "distritos": ["Bagua Grande", "El Milagro", "Utcubamba"]
            },
            "RED BAGUA": {
                "establecimientos": [7225, 7285, 5095, 5096, 7258, 7259, 5066],
                "distritos": ["Bagua", "Aramango", "Copallin"]
            },
            "RED CONDORCANQUI": {
                "establecimientos": [1720, 1744, 1659, 1660, 1661, 1662, 1663, 1664, 1715, 1706, 1681],
                "distritos": ["Santa María de Nieva", "Condorcanqui", "El Cenepa"]
            },
            "RED CHACHAPOYAS": {
                "establecimientos": [2664, 8828, 2570, 1345, 1368, 23961, 3760, 3749],
                "distritos": ["Chachapoyas", "Asunción", "Balsas", "Cheto"]
            }
        }
    
    def generate_presentation(self, filtered_data):
        """Genera una presentación PowerPoint con datos organizados por redes de salud"""
        
        # Crear presentación
        prs = Presentation()
        
        # Diapositiva de título
        self._add_title_slide(prs, filtered_data)
        
        # Diapositiva resumen general
        self._add_summary_slide(prs, filtered_data)
        
        # Diapositivas por red de salud
        for network_name, network_data in self.health_networks.items():
            network_filtered_data = self._filter_data_by_network(filtered_data, network_data["establecimientos"])
            if not network_filtered_data.empty:
                self._add_network_slide(prs, network_name, network_data, network_filtered_data)
                self._add_establishments_detail_slide(prs, network_name, network_data, network_filtered_data)
        
        return prs
    
    def _add_title_slide(self, prs, data):
        """Añade diapositiva de título"""
        slide_layout = prs.slide_layouts[0]  # Layout de título
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Reporte de Vigilancia Epidemiológica"
        subtitle.text = f"Sistema de Vigilancia de Salud\nFecha: {datetime.now().strftime('%d/%m/%Y')}\nTotal de registros: {len(data):,}"
    
    def _add_summary_slide(self, prs, data):
        """Añade diapositiva de resumen general"""
        slide_layout = prs.slide_layouts[5]  # Layout en blanco
        slide = prs.slides.add_slide(slide_layout)
        
        # Título
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "RESUMEN GENERAL DEL SISTEMA"
        title_p.font.size = Pt(24)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.CENTER
        
        # Métricas generales
        metrics = self._calculate_general_metrics(data)
        metrics_text = f"""
📊 ESTADÍSTICAS GENERALES

🏠 Total de viviendas inspeccionadas: {metrics['total_viviendas']:,}
🔍 Total de inspecciones realizadas: {len(data):,}
💧 Consumo total de larvicida: {metrics['consumo_total']:.2f} g
🦟 Total de contenedores tratados: {metrics['contenedores_tratados']:,}
👥 Inspectores activos: {metrics['inspectores_activos']:,}

📈 ÍNDICES DE COBERTURA

🎯 Cobertura promedio: {metrics['cobertura_promedio']:.1f}%
🏆 Eficiencia de tratamiento: {metrics['eficiencia_tratamiento']:.1f}%
        """
        
        content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        content_frame = content_shape.text_frame
        content_p = content_frame.paragraphs[0]
        content_p.text = metrics_text
        content_p.font.size = Pt(14)
    
    def _add_network_slide(self, prs, network_name, network_data, network_filtered_data):
        """Añade diapositiva por red de salud"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Título de la red
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = network_name
        title_p.font.size = Pt(20)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.CENTER
        
        # Métricas de la red
        network_metrics = self._calculate_network_metrics(network_filtered_data)
        
        # Información general
        general_info = f"""
🌐 INFORMACIÓN GENERAL

📍 Distritos: {', '.join(network_data['distritos'])}
🏥 Establecimientos: {len(network_data['establecimientos'])}
🏠 Viviendas inspeccionadas: {network_metrics['total_viviendas']:,}
📊 Total de registros: {len(network_filtered_data):,}
        """
        
        info_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(3))
        info_frame = info_shape.text_frame
        info_p = info_frame.paragraphs[0]
        info_p.text = general_info
        info_p.font.size = Pt(12)
        
        # Métricas operativas
        operational_info = f"""
⚡ MÉTRICAS OPERATIVAS

💧 Consumo de larvicida: {network_metrics['consumo_total']:.2f} g
🦟 Contenedores tratados: {network_metrics['contenedores_tratados']:,}
👥 Inspectores activos: {network_metrics['inspectores_activos']:,}
📈 Cobertura promedio: {network_metrics['cobertura_promedio']:.1f}%
🎯 Índice de positividad: {network_metrics['indice_positividad']:.1f}%
        """
        
        op_shape = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(3))
        op_frame = op_shape.text_frame
        op_p = op_frame.paragraphs[0]
        op_p.text = operational_info
        op_p.font.size = Pt(12)
    
    def _add_establishments_detail_slide(self, prs, network_name, network_data, network_filtered_data):
        """Añade diapositiva detallada de establecimientos"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        # Título
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = f"{network_name} - DETALLE POR ESTABLECIMIENTOS"
        title_p.font.size = Pt(18)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.CENTER
        
        # Crear tabla de establecimientos
        establishment_details = self._get_establishment_details(network_filtered_data, network_data["establecimientos"])
        
        # Añadir detalles como texto
        details_text = "ESTABLECIMIENTO | VIVIENDAS | CONSUMO (g) | COBERTURA\n" + "="*60 + "\n"
        
        for est_id in network_data["establecimientos"]:
            if est_id in self.data_processor.health_facilities:
                est_name = self.data_processor.health_facilities[est_id]["name"]
                est_data = network_filtered_data[network_filtered_data['cod_renipress'] == est_id]
                
                if not est_data.empty:
                    viviendas = len(est_data[est_data['atencion_vivienda_indicador'] == 1])
                    consumo = est_data['consumo_larvicida'].sum()
                    total_houses = self.data_processor.health_facilities[est_id]["total_houses"]
                    cobertura = (viviendas / total_houses * 100) if total_houses > 0 else 0
                    
                    details_text += f"{est_name[:25]:<25} | {viviendas:>8} | {consumo:>10.1f} | {cobertura:>8.1f}%\n"
        
        content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        content_frame = content_shape.text_frame
        content_p = content_frame.paragraphs[0]
        content_p.text = details_text
        content_p.font.size = Pt(10)
        content_p.font.name = "Courier New"  # Monospace para alineación
    
    def _filter_data_by_network(self, data, establishment_ids):
        """Filtra datos por establecimientos de una red"""
        return data[data['cod_renipress'].isin(establishment_ids)]
    
    def _calculate_general_metrics(self, data):
        """Calcula métricas generales"""
        total_viviendas = len(data[data['atencion_vivienda_indicador'] == 1])
        consumo_total = data['consumo_larvicida'].sum()
        
        # Obtener todas las columnas de contenedores organizadas por tipo
        container_types = self.data_processor.get_container_columns()
        
        # Contenedores tratados (suma de todos los tratamientos)
        treated_columns = []
        inspected_columns = []
        positive_columns = []
        
        for container_type_cols in container_types.values():
            for col in container_type_cols:
                if col.endswith('_TQ') or col.endswith('_TF'):
                    treated_columns.append(col)
                elif col.endswith('_I'):
                    inspected_columns.append(col)
                elif col.endswith('_P'):
                    positive_columns.append(col)
        
        contenedores_tratados = 0
        for col in treated_columns:
            if col in data.columns:
                contenedores_tratados += data[col].sum()
        
        # Total inspeccionados
        total_inspeccionados = 0
        for col in inspected_columns:
            if col in data.columns:
                total_inspeccionados += data[col].sum()
        
        # Total positivos
        total_positivos = 0
        for col in positive_columns:
            if col in data.columns:
                total_positivos += data[col].sum()
        
        inspectores_activos = data['usuario_registra'].nunique()
        
        # Cobertura promedio
        cobertura_promedio = 0
        total_establishments = 0
        facility_column = 'cod_renipress'  # Usar columna estándar
        
        for est_id, est_info in self.data_processor.health_facilities.items():
            est_data = data[data[facility_column] == est_id]
            if not est_data.empty:
                viviendas_inspeccionadas = len(est_data[est_data['atencion_vivienda_indicador'] == 1])
                cobertura = (viviendas_inspeccionadas / est_info["total_houses"] * 100) if est_info["total_houses"] > 0 else 0
                cobertura_promedio += cobertura
                total_establishments += 1
        
        cobertura_promedio = cobertura_promedio / total_establishments if total_establishments > 0 else 0
        
        # Eficiencia de tratamiento
        eficiencia_tratamiento = (contenedores_tratados / total_inspeccionados * 100) if total_inspeccionados > 0 else 0
        
        return {
            'total_viviendas': total_viviendas,
            'consumo_total': consumo_total,
            'contenedores_tratados': contenedores_tratados,
            'inspectores_activos': inspectores_activos,
            'cobertura_promedio': cobertura_promedio,
            'eficiencia_tratamiento': eficiencia_tratamiento
        }
    
    def _calculate_network_metrics(self, network_data):
        """Calcula métricas específicas de una red"""
        total_viviendas = len(network_data[network_data['atencion_vivienda_indicador'] == 1])
        consumo_total = network_data['consumo_larvicida'].sum()
        inspectores_activos = network_data['usuario_registra'].nunique()
        
        # Obtener todas las columnas de contenedores organizadas por tipo
        container_types = self.data_processor.get_container_columns()
        
        # Separar columnas por tipo
        treated_columns = []
        inspected_columns = []
        positive_columns = []
        
        for container_type_cols in container_types.values():
            for col in container_type_cols:
                if col.endswith('_TQ') or col.endswith('_TF'):
                    treated_columns.append(col)
                elif col.endswith('_I'):
                    inspected_columns.append(col)
                elif col.endswith('_P'):
                    positive_columns.append(col)
        
        # Contenedores tratados
        contenedores_tratados = 0
        for col in treated_columns:
            if col in network_data.columns:
                contenedores_tratados += network_data[col].sum()
        
        # Total inspeccionados
        total_inspeccionados = 0
        for col in inspected_columns:
            if col in network_data.columns:
                total_inspeccionados += network_data[col].sum()
        
        # Total positivos
        total_positivos = 0
        for col in positive_columns:
            if col in network_data.columns:
                total_positivos += network_data[col].sum()
        
        # Cobertura promedio para la red
        cobertura_promedio = 0
        total_establishments = 0
        
        establishment_ids = network_data['cod_renipress'].unique()
        for est_id in establishment_ids:
            if est_id in self.data_processor.health_facilities:
                est_data = network_data[network_data['cod_renipress'] == est_id]
                viviendas_inspeccionadas = len(est_data[est_data['atencion_vivienda_indicador'] == 1])
                total_houses = self.data_processor.health_facilities[est_id]["total_houses"]
                cobertura = (viviendas_inspeccionadas / total_houses * 100) if total_houses > 0 else 0
                cobertura_promedio += cobertura
                total_establishments += 1
        
        cobertura_promedio = cobertura_promedio / total_establishments if total_establishments > 0 else 0
        
        # Índice de positividad
        indice_positividad = (total_positivos / total_inspeccionados * 100) if total_inspeccionados > 0 else 0
        
        return {
            'total_viviendas': total_viviendas,
            'consumo_total': consumo_total,
            'contenedores_tratados': contenedores_tratados,
            'inspectores_activos': inspectores_activos,
            'cobertura_promedio': cobertura_promedio,
            'indice_positividad': indice_positividad
        }
    
    def _get_establishment_details(self, network_data, establishment_ids):
        """Obtiene detalles de establecimientos"""
        details = []
        for est_id in establishment_ids:
            if est_id in self.data_processor.health_facilities:
                est_data = network_data[network_data['cod_renipress'] == est_id]
                if not est_data.empty:
                    details.append({
                        'id': est_id,
                        'name': self.data_processor.health_facilities[est_id]["name"],
                        'viviendas': len(est_data[est_data['atencion_vivienda_indicador'] == 1]),
                        'consumo': est_data['consumo_larvicida'].sum(),
                        'total_houses': self.data_processor.health_facilities[est_id]["total_houses"]
                    })
        return details
    
    def save_presentation(self, presentation, filename="reporte_vigilancia.pptx"):
        """Guarda la presentación"""
        # Crear directorio de reportes si no existe
        os.makedirs("reportes", exist_ok=True)
        
        # Generar nombre único con timestamp
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_filename = f"reportes/reporte_vigilancia_{timestamp}.pptx"
        
        presentation.save(unique_filename)
        return unique_filename